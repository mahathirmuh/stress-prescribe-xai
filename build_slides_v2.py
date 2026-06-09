# -*- coding: utf-8 -*-
"""
Generate kk-ngebut_slides_v2.pptx — an 18-slide academic deck for the
"Explainable ML + Counterfactual + GenAI" stress-prediction project.

Slide text: English (international conference audience).
Speaker notes: Bahasa Indonesia (for the oral presentation).

Numbers are sourced from outputs/reports/*.csv|json and outputs/recommendations/*.json.
Run:  python build_slides_v2.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ----------------------------------------------------------------------------
# 0. Custom figure — ablation comparison (Locked vs Unlocked)
# ----------------------------------------------------------------------------
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager

FIG_DIR = "outputs/figures"
ABLATION_FIG = os.path.join(FIG_DIR, "ablation_comparison.png")

def make_ablation_figure():
    metrics = ["CF success rate", "% CFs editing\nlocked outcomes", "% behaviour-only\n(actionable) CFs"]
    locked   = [62.5, 0.0, 100.0]     # ours: 25/40 valid, all behaviour-only
    unlocked = [100.0, 100.0, 0.0]    # ablation: 40/40 valid, all touch outcomes
    x = range(len(metrics)); w = 0.36
    fig, ax = plt.subplots(figsize=(7.6, 3.5), dpi=200)
    b1 = ax.bar([i - w/2 for i in x], locked,   w, label="Locked outcomes (ours)",      color="#129E8A")
    b2 = ax.bar([i + w/2 for i in x], unlocked, w, label="Unlocked outcomes (ablation)", color="#C0393D")
    for bars in (b1, b2):
        for b in bars:
            ax.annotate(f"{b.get_height():g}%", (b.get_x()+b.get_width()/2, b.get_height()),
                        ha="center", va="bottom", fontsize=10, fontweight="bold", color="#1E2733",
                        xytext=(0, 2), textcoords="offset points")
    ax.set_ylim(0, 118); ax.set_ylabel("Percentage (%)", fontsize=10, color="#1E2733")
    ax.set_xticks(list(x)); ax.set_xticklabels(metrics, fontsize=9.5, color="#1E2733")
    ax.set_title("Ablation: effect of locking outcome features in DiCE (n = 40 instances)",
                 fontsize=11, fontweight="bold", color="#14234A", pad=10)
    ax.legend(loc="upper center", ncol=2, frameon=False, fontsize=9.5, bbox_to_anchor=(0.5, -0.16))
    ax.spines[["top", "right"]].set_visible(False)
    ax.tick_params(length=0); ax.grid(axis="y", color="#E4E8F0", lw=0.8)
    ax.set_axisbelow(True)
    fig.tight_layout()
    fig.savefig(ABLATION_FIG, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print("wrote", ABLATION_FIG)

make_ablation_figure()

# ----------------------------------------------------------------------------
# 1. Theme
# ----------------------------------------------------------------------------
NAVY   = RGBColor(0x14, 0x23, 0x4A)
INDIGO = RGBColor(0x2B, 0x37, 0x86)
BLUE   = RGBColor(0x2E, 0x6F, 0xE0)   # PREDICT
TEAL   = RGBColor(0x12, 0x9E, 0x8A)   # EXPLAIN / good
AMBER  = RGBColor(0xE0, 0x8A, 0x1E)   # PRESCRIBE / caution
PURPLE = RGBColor(0x7A, 0x4F, 0xC0)   # NATURALIZE
RED    = RGBColor(0xC0, 0x39, 0x3D)   # bad
LIGHT  = RGBColor(0xF5, 0xF7, 0xFB)   # page bg
TINT   = RGBColor(0xEC, 0xF1, 0xFA)   # callout bg
INK    = RGBColor(0x20, 0x28, 0x36)
MUTE   = RGBColor(0x5C, 0x65, 0x72)
LINE   = RGBColor(0xD7, 0xDD, 0xE8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ICE    = RGBColor(0xDD, 0xE6, 0xF6)   # light text on navy

FONT   = "Segoe UI"
FONT_SB= "Segoe UI Semibold"

EMW, EMH = Inches(13.333), Inches(7.5)
PW = 13.333
TOTAL = 17

prs = Presentation()
prs.slide_width  = EMW
prs.slide_height = EMH
BLANK = prs.slide_layouts[6]

# ----------------------------------------------------------------------------
# 2. Low-level helpers
# ----------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)

def _no_line(shape):
    shape.line.fill.background()

def _no_shadow(shape):
    shape.shadow.inherit = False

def fill_rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line_color=None, line_w=1.0, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line_color is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    _no_shadow(sp)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    return sp

def bg(s, color=LIGHT):
    return fill_rect(s, 0, 0, PW, 7.5, color)

def _set_run(r, t, size, color, bold=False, italic=False, font=FONT):
    r.text = t; f = r.font
    f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.italic = italic; f.name = font

def tb(s, x, y, w, h, text, size, color, bold=False, italic=False,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, ls=1.0, shrink=False):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if ls: p.line_spacing = ls
        _set_run(p.add_run(), ln, size, color, bold, italic, font)
    return box

def set_text_in(shape, text, size, color, bold=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE, font=FONT, ls=1.0):
    tf = shape.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.08))
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if ls: p.line_spacing = ls
        _set_run(p.add_run(), ln, size, color, bold, font=font)
    return shape

def bullets(s, x, y, w, h, items, base_size=15, gap=7):
    """items: list of dicts {t, lvl(0/1), color, bold, glyph, gcolor, size}"""
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    for i, it in enumerate(items):
        lvl   = it.get("lvl", 0)
        color = it.get("color", INK)
        bold  = it.get("bold", False)
        size  = it.get("size", base_size if lvl == 0 else base_size - 2)
        glyph = it.get("glyph", "●" if lvl == 0 else "–")
        gcol  = it.get("gcolor", BLUE if lvl == 0 else MUTE)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.line_spacing = it.get("ls", 1.04)
        p.space_after = Pt(it.get("gap", gap))
        pPr = p._p.get_or_add_pPr()
        marL = int(Inches(0.30 + 0.34 * lvl)); indent = int(Inches(-0.30))
        pPr.set("marL", str(marL)); pPr.set("indent", str(indent))
        gsz = 9 if lvl == 0 else 10
        _set_run(p.add_run(), glyph + "  ", gsz if glyph in ("●",) else size, gcol, bold=True)
        _set_run(p.add_run(), it["t"], size, color, bold)
        # optional trailing emphasis run
        if "t2" in it:
            _set_run(p.add_run(), it["t2"], size, it.get("c2", BLUE), bold=True)
    return box

def header(s, title, kicker=None, accent=BLUE):
    bg(s)
    fill_rect(s, 0.55, 0.52, 0.13, 0.66, accent, radius=None)
    ky = 0.50
    if kicker:
        tb(s, 0.86, 0.48, 11.5, 0.3, kicker.upper(), 11.5, accent, bold=True, font=FONT_SB)
        ky = 0.74
    else:
        ky = 0.58
    tb(s, 0.86, ky, 11.9, 0.62, title, 27, NAVY, bold=True, font=FONT_SB)
    fill_rect(s, 0.55, 1.46, 12.23, 0.022, LINE)
    footer(s)

def footer(s):
    fill_rect(s, 0.55, 6.98, 12.23, 0.018, LINE)
    tb(s, 0.55, 7.04, 8.0, 0.3, "Stress Prediction  ·  Explainable ML + Counterfactual + GenAI",
       9, MUTE)
    tb(s, 10.0, 7.04, 2.78, 0.3, f"{s._page} / {TOTAL}", 9, MUTE, align=PP_ALIGN.RIGHT)

def add_image(s, path, x, y, max_w, max_h, center=True, border=True):
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = max_w; h = w / ar
    if h > max_h:
        h = max_h; w = h * ar
    px = x + (max_w - w) / 2 if center else x
    pic = s.shapes.add_picture(path, Inches(px), Inches(y), Inches(w), Inches(h))
    if border:
        pic.line.color.rgb = LINE; pic.line.width = Pt(0.75)
    return pic

def chip(s, x, y, w, h, text, fill, txt=WHITE, size=11, radius=0.5, bold=True):
    sp = fill_rect(s, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)
    set_text_in(sp, text, size, txt, bold=bold, align=PP_ALIGN.CENTER, font=FONT_SB)
    return sp

def callout(s, x, y, w, h, text, accent=BLUE, fill=TINT, size=13, bold_lead=None):
    fill_rect(s, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    fill_rect(s, x, y, 0.09, h, accent)
    box = s.shapes.add_textbox(Inches(x + 0.28), Inches(y), Inches(w - 0.5), Inches(h))
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.line_spacing = 1.06
    if bold_lead:
        _set_run(p.add_run(), bold_lead, size, accent, bold=True, font=FONT_SB)
    _set_run(p.add_run(), text, size, INK)
    return box

def notes(s, text):
    s.notes_slide.notes_text_frame.text = text

# track page numbers
_orig_add = type(prs.slides).add_slide
_page = {"n": 0}
def new_slide():
    s = slide(); _page["n"] += 1; s._page = _page["n"]; return s

# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = new_slide()
bg(s, NAVY)
# accent band
fill_rect(s, 0, 0, PW, 0.16, BLUE)
fill_rect(s, 0, 7.34, PW, 0.16, TEAL)
tb(s, 1.0, 0.95, 11.3, 0.35, "MASTER PROGRAM IN INFORMATICS  ·  ITS SURABAYA",
   12.5, RGBColor(0x9F, 0xB4, 0xDD), bold=True, font=FONT_SB)
tb(s, 1.0, 1.55, 11.3, 2.4,
   "Explainable Machine Learning with\nCounterfactual Analysis and GenAI Naturalization",
   34, WHITE, bold=True, font=FONT_SB, ls=1.05)
tb(s, 1.0, 3.45, 11.3, 0.9,
   "for Stress Prediction and Intervention Using Sleep and Lifestyle Data",
   19, ICE, italic=True, ls=1.05)
fill_rect(s, 1.0, 4.55, 4.2, 0.03, BLUE)
tb(s, 1.0, 4.78, 11.3, 0.55,
   "Obi Kastanya   ·   Dhayu Intan Nareswari   ·   Ananta Dwi Prayoga Alwy   ·   Mahathir Muhammad",
   14.5, WHITE, bold=True)
tb(s, 1.0, 5.28, 11.3, 0.7,
   "Department of Informatics, Faculty of Intelligent Electrical and Informatics Technology\n"
   "Institut Teknologi Sepuluh Nopember (ITS), Surabaya",
   12.5, RGBColor(0xB9, 0xC6, 0xE2), ls=1.1)
# stage chips
cw, gap = 2.55, 0.2
x0 = 1.0
for i, (lbl, col) in enumerate([("1 · PREDICT", BLUE), ("2 · EXPLAIN", TEAL),
                                ("3 · PRESCRIBE", AMBER), ("4 · NATURALIZE", PURPLE)]):
    chip(s, x0 + i * (cw + gap), 6.25, cw, 0.5, lbl, col, size=12)
tb(s, 10.0, 7.04, 2.78, 0.3, f"1 / {TOTAL}", 9, RGBColor(0x8A,0x9A,0xBE), align=PP_ALIGN.RIGHT)
notes(s,
"Selamat pagi/siang/sore Bapak/Ibu juri dan rekan-rekan sekalian. Perkenalkan, kami dari "
"Master Program Informatika, Fakultas Teknologi Elektro dan Informatika Cerdas, Institut "
"Teknologi Sepuluh Nopember Surabaya. Hari ini kami mempresentasikan penelitian berjudul "
"\"Explainable Machine Learning with Counterfactual Analysis and GenAI Naturalization for "
"Stress Prediction and Intervention Using Sleep and Lifestyle Data\". Penelitian ini "
"mengintegrasikan empat tahap: prediksi stres, penjelasan model, rekomendasi intervensi, "
"dan naturalisasi hasil dalam bahasa awam.")

# ============================================================================
# SLIDE 2 — Outline
# ============================================================================
s = new_slide()
header(s, "Outline", "Agenda")
left = [
    ("01", "Background & Motivation", BLUE),
    ("02", "Research Gap & Contributions", BLUE),
    ("03", "Proposed 4-Stage Framework", TEAL),
    ("04", "Dataset & Preprocessing", TEAL),
    ("05", "Stage 1 — Prediction & Models", BLUE),
]
right = [
    ("06", "Stage 2 — SHAP Explainability", TEAL),
    ("07", "Stage 3 — Counterfactuals + Ablation", AMBER),
    ("08", "Stage 4 — GenAI Naturalization", PURPLE),
    ("09", "Case Studies & Limitations", RED),
    ("10", "Conclusion & Summary", RED),
]
def outline_col(items, x):
    y = 1.95
    for num, txt, col in items:
        chip(s, x, y, 0.62, 0.62, num, col, size=15, radius=0.18)
        tb(s, x + 0.85, y + 0.02, 4.9, 0.6, txt, 15.5, INK, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.96
outline_col(left, 0.9)
outline_col(right, 7.0)
notes(s,
"Presentasi ini dibagi ke beberapa bagian. Mulai dari latar belakang dan motivasi, lalu gap "
"penelitian dan kontribusi yang kami tawarkan. Kemudian framework 4-tahap yang kami usulkan "
"beserta dataset dan metodologi setiap tahap. Setelah itu masuk ke hasil — performa model, "
"SHAP, kualitas counterfactual — lalu ablation study untuk membuktikan pentingnya outcome "
"locking, tiga studi kasus individu, keterbatasan, arah ke depan, dan kesimpulan. Estimasi "
"sekitar 15 menit, sisanya untuk diskusi.")

# ============================================================================
# SLIDE 3 — Background & Motivation
# ============================================================================
s = new_slide()
header(s, "Background & Motivation", "Background", BLUE)
bullets(s, 0.9, 1.85, 7.0, 5.0, [
    {"t":"Mental stress is a major modern health concern", "bold":True},
    {"t":"Driven by sleep, lifestyle, work pressure, and digital habits", "lvl":1},
    {"t":"Increasingly trackable via wearables and self-report data", "lvl":1},
    {"t":"ML can predict stress accurately — but most models are black boxes", "bold":True, "gcolor":BLUE},
    {"t":"Accurate yet opaque: they don't explain WHY", "lvl":1},
    {"t":"SHAP / LIME add explanations, but output stays technical", "lvl":1},
    {"t":"What end-users actually need:", "bold":True, "gcolor":BLUE},
    {"t":"PREDICT the stress level", "lvl":1},
    {"t":"EXPLAIN the drivers", "lvl":1},
    {"t":"PRESCRIBE concrete, actionable changes", "lvl":1},
    {"t":"NATURALIZE everything into plain, empathetic language", "lvl":1},
])
# right visual: gap ladder
fill_rect(s, 8.35, 1.95, 4.35, 4.6, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
          radius=0.04, line_color=LINE, line_w=1.0)
tb(s, 8.6, 2.15, 3.9, 0.4, "From prediction to action", 13.5, NAVY, bold=True, font=FONT_SB)
steps = [("Predict", "stress level", BLUE), ("Explain", "feature drivers", TEAL),
         ("Prescribe", "what to change", AMBER), ("Naturalize", "plain narrative", PURPLE)]
yy = 2.75
for i,(a,b,c) in enumerate(steps):
    chip(s, 8.6, yy, 3.85, 0.72, "", c, radius=0.12)
    tb(s, 8.85, yy+0.07, 3.4, 0.3, a, 13.5, WHITE, bold=True, font=FONT_SB)
    tb(s, 8.85, yy+0.37, 3.4, 0.3, b, 11, RGBColor(0xEE,0xF3,0xFF))
    if i < 3:
        tb(s, 10.3, yy+0.70, 0.5, 0.28, "▼", 11, MUTE, align=PP_ALIGN.CENTER)
    yy += 0.94
notes(s,
"Stres mental adalah salah satu masalah kesehatan modern yang paling pelik — dipengaruhi "
"tekanan kerja, pola tidur, aktivitas fisik, kondisi kesehatan, hingga kebiasaan digital "
"seperti screen time sebelum tidur. Dengan data lifestyle dan health, machine learning bisa "
"memprediksi tingkat stres. Tapi kebanyakan model black-box — akurat tapi tidak menjelaskan "
"kenapa. SHAP atau LIME memberi explanation, tapi outputnya tetap teknis, sulit dipahami "
"orang awam. Yang dibutuhkan adalah sistem yang tidak hanya MEMPREDIKSI, tapi juga "
"MENJELASKAN penyebab, MEMBERI SARAN konkret, dan MENATURALISASI semua itu ke bahasa awam.")

# ============================================================================
# SLIDE 4 — Research Gap & Contributions
# ============================================================================
s = new_slide()
header(s, "Research Gap & Contributions", "Background", BLUE)
# Gap card
fill_rect(s, 0.9, 1.9, 5.6, 4.7, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
          radius=0.03, line_color=LINE, line_w=1.0)
fill_rect(s, 0.9, 1.9, 5.6, 0.62, RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
tb(s, 1.1, 1.98, 5.2, 0.45, "Gap in prior work", 15, WHITE, bold=True, font=FONT_SB,
   anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 1.15, 2.78, 5.1, 3.7, [
    {"t":"Stress + SHAP — predictive, not prescriptive", "size":13.5},
    {"t":"No concrete recommendation to act on", "lvl":1, "size":12},
    {"t":"DiCE — common in classification (e.g. loans)", "size":13.5},
    {"t":"Rarely used for regression with causal constraints", "lvl":1, "size":12},
    {"t":"GenAI in health — summarization / education", "size":13.5},
    {"t":"Not used to naturalize counterfactual ML output", "lvl":1, "size":12},
], gap=6)
# Contributions card
fill_rect(s, 6.9, 1.9, 5.55, 4.7, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
          radius=0.03, line_color=LINE, line_w=1.0)
fill_rect(s, 6.9, 1.9, 5.55, 0.62, TEAL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
tb(s, 7.1, 1.98, 5.2, 0.45, "Our contributions", 15, WHITE, bold=True, font=FONT_SB,
   anchor=MSO_ANCHOR.MIDDLE)
bullets(s, 7.15, 2.78, 5.05, 3.7, [
    {"t":"Integrated 4-stage prescriptive framework", "bold":True, "gcolor":TEAL, "size":13.5},
    {"t":"Causally-restricted counterfactuals", "bold":True, "gcolor":TEAL, "size":13.5},
    {"t":"behaviour vs outcome vs immutable features", "lvl":1, "size":12},
    {"t":"GenAI layer with 3-stage validation", "bold":True, "gcolor":TEAL, "size":13.5},
    {"t":"safety + structural + faithfulness checks", "lvl":1, "size":12},
    {"t":"Empirical evidence", "bold":True, "gcolor":TEAL, "size":13.5},
    {"t":"62.5% CF success · 100% plausibility on 40 instances", "lvl":1, "size":12},
], gap=6)
notes(s,
"Di literatur sebelumnya, prediksi stres dan SHAP umumnya sudah ada — tapi belum prescriptive, "
"tidak ada rekomendasi konkret. DiCE counterfactual populer, tapi kebanyakan untuk klasifikasi "
"seperti loan approval; penggunaan untuk regresi dengan constraint kausal strict belum banyak. "
"GenAI di healthcare biasanya untuk summarization atau patient education; belum ada yang "
"memakai GPT sebagai naturaliser output counterfactual. Kontribusi kami ada empat: framework "
"4-tahap terintegrasi; counterfactual yang causally restricted; layer GenAI dengan safety dan "
"faithfulness validation; dan bukti empiris — 62,5 persen CF success dengan 100 persen "
"plausibility pada 40 instance.")

# ============================================================================
# SLIDE 5 — Proposed 4-Stage Framework (diagram)
# ============================================================================
s = new_slide()
header(s, "Proposed 4-Stage Framework", "Framework", TEAL)
tb(s, 0.9, 1.62, 11.5, 0.35, "End-to-end pipeline: from raw sleep/lifestyle data to a narrative recommendation",
   13.5, MUTE, italic=True)
stages = [
    ("1", "PREDICT", BLUE, "CatBoost\nRandom Forest\nTabNet", "stress_score ∈ [1, 10]"),
    ("2", "EXPLAIN", TEAL, "SHAP\nTreeExplainer\nglobal + local", "feature contributions"),
    ("3", "PRESCRIBE", AMBER, "DiCE genetic CF\n+ permitted_range\nbehaviour-only", "behaviour-only CFs"),
    ("4", "NATURALIZE", PURPLE, "GPT-4o-mini\n+ 3-layer\nvalidation", "narrative JSON"),
]
cw = 2.72; gapx = (12.23 - 4*cw) / 3; x = 0.55; cardy = 2.25; ch = 2.55
centers = []
for i,(num,name,col,body,out) in enumerate(stages):
    cx = x + i*(cw+gapx)
    centers.append(cx + cw/2)
    fill_rect(s, cx, cardy, cw, ch, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    tb(s, cx, cardy+0.18, cw, 0.4, f"STAGE {num}", 11, RGBColor(0xEC,0xF1,0xFF), bold=True,
       align=PP_ALIGN.CENTER, font=FONT_SB)
    tb(s, cx, cardy+0.52, cw, 0.45, name, 19, WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_SB)
    fill_rect(s, cx+0.7, cardy+1.04, cw-1.4, 0.018, RGBColor(0xFF,0xFF,0xFF))
    tb(s, cx, cardy+1.18, cw, 1.2, body, 12.5, RGBColor(0xF3,0xF6,0xFF),
       align=PP_ALIGN.CENTER, ls=1.12)
    # output caption
    tb(s, cx, cardy+ch+0.12, cw, 0.6, "▼", 12, col, align=PP_ALIGN.CENTER)
    tb(s, cx, cardy+ch+0.40, cw, 0.5, out, 11.5, INK, align=PP_ALIGN.CENTER, bold=True)
    # arrow to next
    if i < 3:
        ar = fill_rect(s, cx+cw+0.06, cardy+ch/2-0.18, gapx-0.12, 0.36, RGBColor(0xB7,0xC0,0xCE),
                       shape=MSO_SHAPE.RIGHT_ARROW)
callout(s, 0.55, 5.78, 12.23, 1.0,
        "only BEHAVIOUR features may change (sleep duration, screen time, caffeine, exercise, …). "
        "OUTCOME and IMMUTABLE features are LOCKED so every recommendation is a manipulable cause — not an outcome the user cannot control.",
        accent=AMBER, fill=RGBColor(0xFD,0xF4,0xE6), bold_lead="Causal restriction:  ")
notes(s,
"Ini gambaran umum framework, mengalir dari kiri ke kanan. Tahap pertama PREDICT: tiga model "
"regresi — CatBoost, Random Forest, TabNet — memprediksi stress_score 1.0 sampai 10.0. Tahap "
"kedua EXPLAIN: SHAP TreeExplainer memberi penjelasan global dan lokal. Tahap ketiga PRESCRIBE: "
"DiCE meng-generate counterfactual, yaitu perubahan minimal yang menurunkan prediksi stres — "
"dan hanya fitur BEHAVIOR yang boleh diubah; outcome dan immutable di-lock. Tahap keempat "
"NATURALIZE: GPT-4o-mini mengubah angka teknis menjadi narasi JSON terstruktur Bahasa "
"Indonesia, dengan safety filter dan faithfulness validation. Catatan di bawah penting: "
"restriksi causal inilah yang menjamin rekomendasi kami benar-benar actionable.")

# ============================================================================
# SLIDE 6 — Dataset & Preprocessing
# ============================================================================
s = new_slide()
header(s, "Dataset & Preprocessing", "Dataset", TEAL)
bullets(s, 0.9, 1.85, 7.15, 5.0, [
    {"t":"Sleep Health & Daily Performance Dataset (Kaggle)", "bold":True, "gcolor":TEAL},
    {"t":"100,000 synthetic samples × 32 features", "lvl":1},
    {"t":"Target: stress_score ∈ [1.0, 10.0]  (mean 5.73, SD 1.62)", "lvl":1},
    {"t":"Working subset: 10,000 stratified by stress_score deciles", "lvl":1},
    {"t":"Strongest negative correlation: sleep_quality_score (r = −0.639)", "lvl":1},
    {"t":"Dropped 4 leakage features", "bold":True, "gcolor":TEAL},
    {"t":"person_id · cognitive_performance_score · sleep_disorder_risk · felt_rested", "lvl":1, "size":12},
    {"t":"Dual encoding strategy", "bold":True, "gcolor":TEAL},
    {"t":"A: native categoricals (CatBoost)", "lvl":1},
    {"t":"B: ordinal encoding + StandardScaler (RF, TabNet)", "lvl":1},
    {"t":"Stratified 70 / 15 / 15 train / val / test split", "bold":True, "gcolor":TEAL},
])
add_image(s, os.path.join(FIG_DIR, "eda_stress_distribution.png"), 8.25, 2.4, 4.45, 3.4)
tb(s, 8.25, 5.75, 4.45, 0.4, "Distribution of stress_score in the 10k sample",
   11, MUTE, italic=True, align=PP_ALIGN.CENTER)
notes(s,
"Dataset kami: Sleep Health & Daily Performance dari Kaggle. 100 ribu sampel sintetis, 32 "
"fitur, mencakup pola tidur, lifestyle, kesehatan, dan aktivitas harian. Untuk iterasi kami "
"pakai stratified sample 10 ribu berdasarkan desil stress_score. EDA mengkonfirmasi rata-rata "
"stress_score 5,73 dengan SD 1,62, dan korelasi negatif terkuat adalah sleep_quality_score "
"r = -0,639 — sesuai ekspektasi domain. Kami drop 4 fitur leakage: person_id, "
"cognitive_performance_score, sleep_disorder_risk, dan felt_rested. Kami siapkan dua versi "
"data — Versi A native kategorikal untuk CatBoost, Versi B ordinal + StandardScaler untuk RF "
"dan TabNet. Split stratified 70/15/15.")

# ============================================================================
# SLIDE 7 — Stage 1: Models & Comparison
# ============================================================================
s = new_slide()
header(s, "Stage 1 — Prediction: Models & Comparison", "Stage 1 · Predict", BLUE)
add_image(s, os.path.join(FIG_DIR, "model_comparison.png"), 0.9, 1.75, 11.5, 2.55)
# results table-ish row of chips
data = [("CatBoost ★", "0.6503", "0.9523", "0.7584", BLUE),
        ("TabNet",     "0.6422", "0.9634", "0.7672", MUTE),
        ("RandomForest","0.6239","0.9876", "0.7856", MUTE)]
ty = 4.55
tb(s, 0.9, ty, 3.0, 0.3, "Model", 12, MUTE, bold=True)
tb(s, 4.2, ty, 1.6, 0.3, "R²", 12, MUTE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 5.9, ty, 1.6, 0.3, "RMSE", 12, MUTE, bold=True, align=PP_ALIGN.CENTER)
tb(s, 7.6, ty, 1.6, 0.3, "MAE", 12, MUTE, bold=True, align=PP_ALIGN.CENTER)
fill_rect(s, 0.9, ty+0.32, 8.3, 0.015, LINE)
yy = ty + 0.42
for name, r2, rmse, mae, col in data:
    bold = name.endswith("★")
    tb(s, 0.9, yy, 3.2, 0.3, name, 13, col, bold=bold)
    tb(s, 4.2, yy, 1.6, 0.3, r2, 13, col, bold=bold, align=PP_ALIGN.CENTER)
    tb(s, 5.9, yy, 1.6, 0.3, rmse, 13, col, align=PP_ALIGN.CENTER)
    tb(s, 7.6, yy, 1.6, 0.3, mae, 13, col, align=PP_ALIGN.CENTER)
    yy += 0.40
# right highlight
fill_rect(s, 9.55, 4.5, 3.18, 2.12, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
tb(s, 9.75, 4.66, 2.8, 0.3, "Best model — CatBoost", 12.5, RGBColor(0x9F,0xB4,0xDD), bold=True, font=FONT_SB)
tb(s, 9.75, 5.0, 2.8, 0.5, "R² = 0.6503", 20, WHITE, bold=True, font=FONT_SB)
tb(s, 9.75, 5.5, 2.9, 0.9,
   "5-fold CV R² = 0.6266 ± 0.0093\n→ stable, no overfitting\n1000 iters · depth 6 · early stop",
   11.5, ICE, ls=1.15)
notes(s,
"Kami train tiga model regresi: CatBoost 1000 iterasi depth 6 dengan early stopping, Random "
"Forest 300 estimators, dan TabNet sampai 200 epoch. Hasilnya CatBoost menang di semua metric "
"— R² 0,6503, RMSE 0,9523, MAE 0,7584. TabNet kedua dengan R² 0,6422, tipis di belakang; "
"Random Forest ketiga di 0,6239. Menariknya, gap dua teratas sangat kecil — gradient boosting "
"dan deep learning attention-based sama-sama viable untuk regresi tabular ini. 5-Fold CV pada "
"CatBoost: R² 0,6266 plus minus 0,0093 — variansi kecil, model stabil. Untuk tahap berikutnya "
"kami pakai CatBoost karena performa terbaik plus native categorical support.")

# ============================================================================
# SLIDE 8 — Stage 2: SHAP global + domain validity
# ============================================================================
s = new_slide()
header(s, "Stage 2 — SHAP Explainability (Global)", "Stage 2 · Explain", TEAL)
add_image(s, os.path.join(FIG_DIR, "shap_global_bar.png"), 8.0, 1.8, 4.7, 4.7)
bullets(s, 0.9, 1.85, 6.9, 2.6, [
    {"t":"TreeExplainer on CatBoost (2,000-sample test subset)", "gcolor":TEAL},
    {"t":"Top features by mean |SHAP|:", "bold":True, "gcolor":TEAL},
    {"t":"sleep_quality_score  0.660    ·    occupation  0.569", "lvl":1, "size":13},
    {"t":"sleep_duration_hrs  0.133    ·    room_temperature  0.084", "lvl":1, "size":13},
    {"t":"wake_episodes_per_night  0.077", "lvl":1, "size":13},
], gap=6)
# domain validity panel
fill_rect(s, 0.9, 4.45, 6.9, 2.15, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
          radius=0.04, line_color=TEAL, line_w=1.25)
tb(s, 1.1, 4.58, 6.5, 0.35, "Domain validity sign-check  (corr of feature value vs its SHAP)",
   12.5, TEAL, bold=True, font=FONT_SB)
dv = [("sleep_quality_score", "−0.996"), ("wake_episodes", "−0.981"),
      ("alcohol_units", "−0.977"), ("deep_sleep_%", "−0.963"),
      ("heart_rate_resting", "+0.974")]
xx = 1.1; cwid = 1.32
for i,(f,v) in enumerate(dv):
    cx = 1.1 + i*1.33
    tb(s, cx, 5.05, 1.3, 0.6, f, 9.5, INK, align=PP_ALIGN.CENTER, ls=0.95)
    tb(s, cx, 5.62, 1.3, 0.35, v, 14, TEAL if v.startswith("−") else AMBER, bold=True, align=PP_ALIGN.CENTER)
tb(s, 1.1, 6.08, 6.5, 0.45,
   "Signs match sleep-medicine literature → the model learned valid, not spurious, patterns.",
   11, MUTE, italic=True)
notes(s,
"SHAP TreeExplainer kami aplikasikan pada CatBoost dengan 2000 sample test. Top fitur global: "
"sleep_quality_score paling dominan dengan mean SHAP 0,66, lalu occupation 0,57, kemudian "
"sleep_duration 0,13. Yang penting adalah domain validity sign-check: untuk tiap fitur numerik "
"utama kami hitung korelasi antara nilai fitur dengan kontribusi SHAP-nya. Hasilnya "
"meyakinkan — sleep_quality_score korelasi minus 0,996, hampir sempurna: makin tinggi kualitas "
"tidur, makin rendah kontribusi ke stres. Wake_episodes minus 0,981, alcohol minus 0,977, "
"deep_sleep minus 0,963, dan resting heart rate positif 0,974 — sesuai bahwa HR istirahat "
"tinggi indikator stres. Ini membuktikan model belajar pola yang konsisten dengan literatur "
"sleep medicine, bukan pola spurious.")

# ============================================================================
# SLIDE 9 — Stage 2: SHAP local (waterfall)
# ============================================================================
s = new_slide()
header(s, "Stage 2 — SHAP Local Explanations", "Stage 2 · Explain", TEAL)
tb(s, 0.9, 1.62, 11.5, 0.35,
   "Per-individual waterfall plots decompose each prediction into feature contributions",
   13.5, MUTE, italic=True)
wf = [("shap_waterfall_low.png", "LOW  ·  pred 4.05", BLUE),
      ("shap_waterfall_mid.png", "MID  ·  pred 5.25", AMBER),
      ("shap_waterfall_high.png","HIGH  ·  pred 7.21", RED)]
x0 = 0.7; iw = 4.0; gp = 0.18
for i,(f,lab,col) in enumerate(wf):
    cx = x0 + i*(iw+gp)
    add_image(s, os.path.join(FIG_DIR, f), cx, 2.15, iw, 3.55, border=True)
    chip(s, cx+iw/2-1.1, 5.8, 2.2, 0.46, lab, col, size=12)
callout(s, 0.55, 6.42, 12.23, 0.5,
        "Each bar shows how a feature pushes the prediction above or below the baseline — the foundation for the counterfactual stage.",
        accent=TEAL, size=12)
notes(s,
"Selain global, SHAP juga kami pakai lokal — per individu — lewat waterfall plot. Untuk tiga "
"individu representatif (low, mid, high stress), plot ini memecah prediksi menjadi kontribusi "
"tiap fitur: bar merah mendorong prediksi naik, bar biru menurunkan. Misalnya pada individu "
"HIGH, kualitas tidur rendah dan durasi tidur kurang mendorong skor stres naik. Penjelasan "
"lokal inilah yang menjadi dasar tahap counterfactual — kita tahu fitur mana yang paling "
"berkontribusi sebelum mencari perubahan minimal untuk menurunkannya.")

# ============================================================================
# SLIDE 10 — Stage 3: DiCE Counterfactual Setup
# ============================================================================
s = new_slide()
header(s, "Stage 3 — DiCE Counterfactual Setup", "Stage 3 · Prescribe", AMBER)
tb(s, 0.9, 1.6, 11.5, 0.35,
   "DiCE genetic algorithm, model_type='regressor' — feature roles enforce causal validity",
   13.5, MUTE, italic=True)
cats = [
    ("BEHAVIOUR", TEAL, "VARYABLE",
     "sleep_duration · screen_time · caffeine · alcohol · exercise · steps · nap · room_temp · sleep_aid"),
    ("OUTCOME", AMBER, "LOCKED (symptoms)",
     "sleep_quality_score · rem_% · wake_episodes · sleep_latency · deep_sleep_% · resting_HR"),
    ("IMMUTABLE", RED, "LOCKED (static)",
     "age · gender · occupation · BMI · country · chronotype · season …"),
]
yy = 2.2
for name, col, tag, feats in cats:
    fill_rect(s, 0.9, yy, 11.55, 1.05, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
              radius=0.06, line_color=LINE, line_w=1.0)
    fill_rect(s, 0.9, yy, 2.55, 1.05, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    tb(s, 0.9, yy+0.22, 2.55, 0.4, name, 15, WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_SB)
    tb(s, 0.9, yy+0.6, 2.55, 0.3, tag, 10.5, RGBColor(0xF3,0xF6,0xFF), align=PP_ALIGN.CENTER)
    tb(s, 3.7, yy+0.18, 8.5, 0.7, feats, 12, INK, anchor=MSO_ANCHOR.MIDDLE, ls=1.05)
    yy += 1.18
callout(s, 0.9, 5.85, 11.55, 0.95,
        "alcohol locked at [0, 0] (never recommend alcohol) · sleep_duration ∈ [4, 10] h · caffeine ≤ 400 mg (FDA). "
        "5 metrics evaluated: Validity · Proximity · Sparsity · Diversity · Plausibility.",
        accent=AMBER, fill=RGBColor(0xFD,0xF4,0xE6), bold_lead="permitted_range safety:  ", size=12)
notes(s,
"DiCE kami konfigurasi dengan genetic algorithm karena lebih robust dari metode random untuk "
"regresi dengan desired range. Fitur dikelompokkan jadi tiga kategori penting untuk causal "
"validity. BEHAVIOR — boleh diubah — fitur yang bisa dimanipulasi langsung pengguna: durasi "
"tidur, screen time, kafein, alkohol, exercise, langkah, nap, suhu ruangan, sleep aid. OUTCOME "
"— di-lock — fitur fisiologis yang merupakan gejala, bukan kausa: sleep quality, REM, wake "
"episodes, latency, deep sleep, resting HR. IMMUTABLE — atribut statis seperti usia, gender, "
"occupation, BMI. Permitted_range eksplisit per fitur: alcohol dikunci di [0,0] supaya DiCE "
"tidak pernah menyarankan alkohol, durasi tidur dibatasi 4–10 jam, kafein maksimal 400mg "
"sesuai FDA. Kami evaluasi 5 metrik: validity, proximity, sparsity, diversity, plausibility.")

# ============================================================================
# SLIDE 11 — Stage 3: Counterfactual Quality Results
# ============================================================================
s = new_slide()
header(s, "Counterfactual Quality Results", "Stage 3 · Prescribe", AMBER)
add_image(s, os.path.join(FIG_DIR, "cf_metrics_by_quartile.png"), 7.35, 1.9, 5.35, 4.4)
# metric chips on left
metrics = [("62.5%", "CF success rate (25/40)", AMBER),
           ("100%", "Validity — all reach target", TEAL),
           ("100%", "Plausibility — within permitted_range", TEAL),
           ("0.226", "Mean proximity (close to original)", BLUE),
           ("4.65", "Mean sparsity (features changed)", BLUE)]
yy = 1.95
for val, lab, col in metrics:
    fill_rect(s, 0.9, yy, 6.1, 0.74, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
              radius=0.1, line_color=LINE, line_w=1.0)
    fill_rect(s, 0.9, yy, 1.55, 0.74, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    tb(s, 0.9, yy+0.16, 1.55, 0.45, val, 17, WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_SB)
    tb(s, 2.65, yy+0.04, 4.2, 0.66, lab, 12.5, INK, anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.82
callout(s, 0.9, 6.05, 6.1, 0.82,
        "62.5% < the 80% target is a CONSEQUENCE of locking outcomes — fewer levers, but every CF is honest & actionable. Quartile balance: 7/6/6/6.",
        accent=AMBER, fill=RGBColor(0xFD,0xF4,0xE6), size=11.5)
notes(s,
"Evaluasi CF pada 40 test instance, 10 per quartile stres. Dari 40, 25 berhasil menghasilkan "
"CF valid — 62,5 persen success rate. Yang berhasil semuanya 100 persen valid (mencapai target "
"range) dan 100 persen plausible (dalam permitted_range). Mean proximity 0,226 — relatif kecil, "
"CF tidak jauh dari kondisi awal. Sparsity 4,65 — sekitar 4–5 fitur diubah per CF, masih masuk "
"akal. Success seimbang antar quartile: 7, 6, 6, 6, tidak bias ke level stres tertentu. Yang "
"harus saya disclose: 62,5 persen memang di bawah target awal 80 persen — tapi ini bukan "
"kegagalan, melainkan konsekuensi dari restriksi causal. Karena outcome di-lock, model hanya "
"punya behavior features untuk digeser, yang impact-nya lebih kecil. Ini trade-off causal "
"validity versus success rate.")

# ============================================================================
# SLIDE 12 — Ablation Study: Outcome Locking (NEW)
# ============================================================================
s = new_slide()
header(s, "Ablation Study — Outcome Locking", "Stage 3 · Key finding", PURPLE)
bullets(s, 0.9, 1.85, 5.7, 3.6, [
    {"t":"Question: is the causal restriction actually necessary?", "bold":True, "gcolor":PURPLE},
    {"t":"Re-ran DiCE on the same 40 instances, two setups:", "gcolor":PURPLE},
    {"t":"Setup 1 (ours): vary BEHAVIOUR only → outcomes locked", "lvl":1, "size":12.5},
    {"t":"Setup 2 (ablation): vary BEHAVIOUR + OUTCOME", "lvl":1, "size":12.5},
    {"t":"Result when outcomes are unlocked:", "bold":True, "gcolor":PURPLE},
    {"t":"success jumps to 100% (40/40) …", "lvl":1, "size":12.5},
    {"t":"… but 100% of CFs edit locked outcomes", "lvl":1, "size":12.5, "color":RED, "gcolor":RED},
    {"t":"0% rely on behaviour-only (actionable) changes", "lvl":1, "size":12.5, "color":RED, "gcolor":RED},
], gap=6)
add_image(s, ABLATION_FIG, 6.7, 1.85, 6.0, 3.7, border=True)
callout(s, 0.55, 5.95, 12.23, 0.95,
        "without locking, every counterfactual hits the target by editing physiological outcomes the user CANNOT directly control. "
        "Locking trades 37.5 points of raw success for 100% actionable, causally-honest advice — validating the restriction empirically.",
        accent=PURPLE, fill=RGBColor(0xF3,0xEE,0xFB), bold_lead="Takeaway:  ", size=12.5)
notes(s,
"Ini temuan kunci yang membedakan paper kami. Pertanyaannya: apakah restriksi causal — mengunci "
"outcome — memang perlu? Kami jalankan ulang DiCE pada 40 instance yang sama dengan dua setup. "
"Setup 1, punya kami: hanya behavior yang boleh berubah, outcome dikunci. Setup 2, ablation: "
"behavior plus outcome boleh berubah. Hasilnya dramatis: saat outcome dibuka, success rate "
"melonjak dari 62,5 ke 100 persen — terlihat lebih baik. TAPI, seratus persen counterfactual "
"itu mencapai target dengan mengubah outcome features — misalnya menaikkan langsung "
"sleep_quality_score atau menurunkan resting heart rate — yang justru TIDAK bisa dikontrol "
"langsung oleh pengguna. Nol persen yang murni behavior. Artinya: tanpa locking, DiCE "
"'curang' — ia memberi rekomendasi yang tidak actionable. Locking menukar 37,5 poin raw "
"success demi 100 persen saran yang actionable dan causally honest. Inilah bukti empiris bahwa "
"restriksi causal kami memang diperlukan, bukan sekadar pilihan desain.")

# ============================================================================
# SLIDE 13 — Stage 4: GenAI Naturalization
# ============================================================================
s = new_slide()
header(s, "Stage 4 — GenAI Naturalization", "Stage 4 · Naturalize", PURPLE)
# input -> GPT -> output flow
fill_rect(s, 0.9, 1.95, 3.3, 1.5, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06,
          line_color=LINE, line_w=1.0)
tb(s, 1.05, 2.08, 3.0, 0.35, "INPUT", 12, BLUE, bold=True, font=FONT_SB)
tb(s, 1.05, 2.42, 3.05, 1.0, "SHAP top-5 features\n+ counterfactual changes\n+ user profile",
   11.5, INK, ls=1.12)
fill_rect(s, 4.32, 2.5, 0.55, 0.4, RGBColor(0xB7,0xC0,0xCE), shape=MSO_SHAPE.RIGHT_ARROW)
fill_rect(s, 5.0, 1.95, 3.3, 1.5, PURPLE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
tb(s, 5.15, 2.08, 3.0, 0.35, "GPT-4o-mini", 13.5, WHITE, bold=True, font=FONT_SB)
tb(s, 5.15, 2.45, 3.05, 1.0, "temperature = 0.3\nresponse_format = JSON\nmax 3 retries",
   11.5, RGBColor(0xF3,0xF0,0xFC), ls=1.12)
fill_rect(s, 8.42, 2.5, 0.55, 0.4, RGBColor(0xB7,0xC0,0xCE), shape=MSO_SHAPE.RIGHT_ARROW)
fill_rect(s, 9.1, 1.95, 3.35, 1.5, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06,
          line_color=LINE, line_w=1.0)
tb(s, 9.25, 2.08, 3.0, 0.35, "OUTPUT (JSON)", 12, TEAL, bold=True, font=FONT_SB)
tb(s, 9.25, 2.42, 3.1, 1.0, "summary · drivers ·\nrecommendations ·\nencouragement · disclaimer",
   11.5, INK, ls=1.12)
# 3-layer validation
tb(s, 0.9, 3.75, 11.5, 0.35, "Post-generation: 3-layer validation with iterative refinement",
   14, NAVY, bold=True, font=FONT_SB)
layers = [
    ("Layer 1 — Safety", RED, "regex blocklist: no drugs, no diagnosis,\nno guarantees, never promote alcohol"),
    ("Layer 2 — Faithfulness", AMBER, "before/after numbers must match exactly;\nno locked-outcome feature may appear"),
    ("Layer 3 — Structural", TEAL, "required keys + type + nested schema;\nretry with feedback if invalid"),
]
x0 = 0.9; lw = 3.78; gp = 0.1
for i,(name, col, body) in enumerate(layers):
    cx = x0 + i*(lw+gp)
    fill_rect(s, cx, 4.25, lw, 1.95, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
              radius=0.05, line_color=LINE, line_w=1.0)
    fill_rect(s, cx, 4.25, lw, 0.55, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    tb(s, cx, 4.33, lw, 0.4, name, 13.5, WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_SB)
    tb(s, cx+0.22, 4.95, lw-0.4, 1.15, body, 12, INK, ls=1.15)
callout(s, 0.9, 6.35, 11.55, 0.5,
        "Layer 2 (faithfulness) is the novel safeguard — it guarantees the narrative never invents or rounds the counterfactual numbers.",
        accent=PURPLE, fill=RGBColor(0xF3,0xEE,0xFB), size=12)
notes(s,
"Tahap terakhir adalah naturalisasi pakai GPT-4o-mini. Input GPT: top-5 fitur SHAP, perubahan "
"counterfactual, dan profil pengguna. Output: JSON terstruktur dengan lima field — summary, "
"drivers, recommendations dengan action/target/rationale, encouragement, dan disclaimer. API "
"call dengan temperature 0,3 untuk determinism dan response_format JSON untuk parsing reliable. "
"Yang krusial: tiga layer validasi pasca-generasi dengan maksimal tiga retry. Layer 1 safety "
"filter regex — blokir kata obat, diagnosa, menjamin, promosi alkohol. Layer 2 faithfulness "
"check — ini yang novel: recommendations harus menyebut angka EXACT dari counterfactual dan "
"tidak boleh menyarankan perubahan pada locked outcome. Layer 3 structural validation — semua "
"required key harus ada dengan tipe benar. Output yang gagal di-retry dengan feedback hingga "
"tiga kali sebelum fallback.")

# ============================================================================
# SLIDE 14 — Case Study Narratives
# ============================================================================
s = new_slide()
header(s, "Case Study Narratives", "Results · Case studies", PURPLE)
tb(s, 0.9, 1.6, 11.5, 0.35, "Three individuals across stress levels — all produced valid, actionable counterfactuals",
   13.5, MUTE, italic=True)
people = [("LOW", "4.05", "3.89", "−0.16", BLUE),
          ("MID", "5.25", "4.94", "−0.32", AMBER),
          ("HIGH", "7.21", "6.86", "−0.35", RED)]
yy = 2.1
for name, pred, cf, d, col in people:
    fill_rect(s, 0.9, yy, 5.5, 0.86, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
              radius=0.08, line_color=LINE, line_w=1.0)
    fill_rect(s, 0.9, yy, 1.35, 0.86, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    tb(s, 0.9, yy+0.27, 1.35, 0.4, name, 15, WHITE, bold=True, align=PP_ALIGN.CENTER, font=FONT_SB)
    tb(s, 2.45, yy+0.12, 3.8, 0.32, f"pred {pred}  →  CF {cf}", 13.5, INK, bold=True)
    tb(s, 2.45, yy+0.46, 3.8, 0.3, f"Δ stress = {d}   ·   all CFs successful", 11.5, MUTE)
    yy += 0.98
# GenAI output card (HIGH)
fill_rect(s, 6.7, 2.1, 5.75, 4.45, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
          radius=0.03, line_color=PURPLE, line_w=1.25)
fill_rect(s, 6.7, 2.1, 5.75, 0.55, PURPLE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
tb(s, 6.9, 2.18, 5.4, 0.4, "GenAI narrative — HIGH stress (excerpt)", 13, WHITE, bold=True,
   font=FONT_SB, anchor=MSO_ANCHOR.MIDDLE)
tb(s, 6.95, 2.78, 5.3, 0.85,
   "“Saat ini Anda memiliki tingkat stres yang cukup tinggi dengan skor 7.21 …”",
   12, INK, italic=True, ls=1.12)
recs = [
    "Perpanjang durasi tidur — dari 6.01 → 8.38 jam",
    "Persingkat waktu layar — dari 102 → 98 menit",
    "Hilangkan alkohol — dari 1.0 → 0.0 unit",
    "Tingkatkan suhu ruang — dari 20.8 → 24.5 °C",
]
bullets(s, 6.95, 3.7, 5.35, 2.2, [{"t":r, "size":12.5, "gcolor":PURPLE} for r in recs], gap=7)
tb(s, 6.95, 6.05, 5.4, 0.45,
   "All actions PASS faithfulness — exact numbers, correct direction, no locked feature.",
   11, TEAL, italic=True, bold=True)
notes(s,
"Tiga individu dipilih sebagai case study berdasarkan stress_score aktual: low sekitar 3, mid "
"6, high 8,5. Ketiganya berhasil mendapat CF dengan delta meaningful — sekitar minus 0,16 untuk "
"low, minus 0,32 mid, minus 0,35 high. Contoh output GenAI untuk kasus HIGH sangat menarik. "
"Action 1: perpanjang durasi tidur dari 6,01 ke 8,38 jam — valid secara klinis. Action 2: "
"persingkat waktu layar 102 ke 98 menit. Action 3: hilangkan alkohol dari 1 unit ke 0 — "
"intervensi yang sangat valid. Action 4: naikkan suhu ruangan untuk kenyamanan tidur. Beberapa "
"arah perubahan bisa terlihat counter-intuitif karena murni mengikuti output DiCE pada data "
"sintetis — kami laporkan apa adanya. Yang penting: SEMUA rekomendasi lolos faithfulness "
"validation — angka exact match dengan CF, arah verb benar, dan tidak ada fitur locked yang "
"muncul sebagai action.")

# ============================================================================
# SLIDE 15 — Limitations & Threats to Validity
# ============================================================================
s = new_slide()
header(s, "Limitations & Threats to Validity", "Limitations", RED)
items = [
    ("Synthetic dataset", "Kaggle-generated; correlations may be artifacts. Proof-of-concept, not clinical validation."),
    ("Causal stationarity", "DiCE assumes feature changes map consistently to predictions — unverified on synthetic data."),
    ("GenAI evaluation scope", "Structural + safety + basic faithfulness only; no user study, expert panel, or NLI scoring."),
    ("Reproducibility bound", "DiCE genetic is stochastic; GPT probabilistic even at T=0.3; TabNet seed-sensitive."),
    ("Modest effect size", "Behaviour-only CFs shift ~0.3 pts (direct effect); outcomes mediate the rest (indirect)."),
    ("GPT hallucination risk", "Regex mitigates, but formal NLI-based faithfulness scoring is recommended for publication."),
]
x0 = 0.9; cwid = 5.78; gpx = 0.1; gpy = 0.18; rh = 1.42
for i,(t,b) in enumerate(items):
    col = i % 2; row = i // 2
    cx = x0 + col*(cwid+gpx); cy = 1.95 + row*(rh+gpy)
    fill_rect(s, cx, cy, cwid, rh, WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
              radius=0.05, line_color=LINE, line_w=1.0)
    fill_rect(s, cx, cy, 0.1, rh, RED)
    tb(s, cx+0.3, cy+0.14, cwid-0.5, 0.35, f"{i+1}.  {t}", 13.5, NAVY, bold=True, font=FONT_SB)
    tb(s, cx+0.3, cy+0.55, cwid-0.55, 0.85, b, 11.8, INK, ls=1.08)
notes(s,
"Kami jujur soal keterbatasan. Pertama dan paling penting: dataset SINTETIS, bukan dari subjek "
"manusia nyata — korelasi bisa jadi artifact generator; ini proof-of-concept metodologis, bukan "
"validasi klinis. Kedua: DiCE mengasumsikan causal stationarity, yang belum terjamin pada data "
"sintetis; rekomendasi dipahami sebagai kemungkinan, bukan kepastian. Ketiga: evaluasi GenAI "
"terbatas — structural, safety regex, dan basic faithfulness; tidak ada user study atau NLI "
"scoring. Keempat: reproducibility ada bound-nya — DiCE genetic stokastik, GPT probabilistik "
"meski temperature di-fix. Kelima: effect size modest, sekitar 0,3 poin per CF, karena ini "
"direct effect saja; total effect dimediasi oleh outcomes. Keenam: risiko halusinasi GPT — "
"regex membantu, tapi untuk publikasi sebaiknya ditambah faithfulness scoring berbasis NLI.")

# ============================================================================
# SLIDE 16 — Conclusion
# ============================================================================
s = new_slide()
header(s, "Conclusion", "Conclusion", TEAL)
bullets(s, 0.9, 1.9, 6.6, 4.6, [
    {"t":"An end-to-end framework integrating four stages:", "bold":True, "gcolor":TEAL},
    {"t":"Prediction → Explanation → Prescription → Naturalization", "lvl":1, "size":13},
    {"t":"Empirical evidence:", "bold":True, "gcolor":TEAL},
    {"t":"CatBoost R² 0.6503 · MAE 0.7584 · CV 0.6266 ± 0.0093", "lvl":1, "size":12.5},
    {"t":"SHAP domain validity up to r = −0.996", "lvl":1, "size":12.5},
    {"t":"CF success 62.5% · 100% plausibility · ablation-validated locking", "lvl":1, "size":12.5},
    {"t":"GenAI: 3/3 case studies pass safety + faithfulness", "lvl":1, "size":12.5},
    {"t":"A complete prescriptive ML pipeline — causally aware & safety-validated", "bold":True, "gcolor":TEAL},
    {"t":"Ready for clinical / wearable / longitudinal validation", "lvl":1, "size":12.5},
], gap=7)
# right summary chips
sumk = [("R² 0.6503", "CatBoost best model", BLUE),
        ("−0.996", "SHAP domain validity", TEAL),
        ("62.5%", "CF success (100% plausible)", AMBER),
        ("3/3", "case studies validated", PURPLE)]
yy = 1.95
for v, l, c in sumk:
    fill_rect(s, 7.9, yy, 4.55, 1.05, c, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07)
    tb(s, 8.15, yy+0.14, 4.1, 0.5, v, 22, WHITE, bold=True, font=FONT_SB)
    tb(s, 8.15, yy+0.64, 4.1, 0.32, l, 12, RGBColor(0xF1,0xF5,0xFF))
    yy += 1.18
notes(s,
"Kesimpulan. Kami berhasil membangun framework end-to-end yang mengintegrasikan empat tahap: "
"prediksi, explanation, prescription, dan naturalization. Bukti empiris: CatBoost R² 0,6503, "
"MAE 0,7584, 5-Fold CV 0,6266 dengan variansi sangat kecil; SHAP domain validity teruji hingga "
"korelasi minus 0,996; counterfactual success 62,5 persen dengan plausibility 100 persen, dan "
"pentingnya outcome locking sudah dibuktikan lewat ablation study; GenAI naturalisasi sukses "
"untuk ketiga case study dengan safety dan faithfulness pass. Kontribusi utama: pipeline "
"prescriptive ML yang lengkap, causally aware, dan safety validated — siap untuk validasi "
"lanjutan di setting klinis, wearable, atau longitudinal.")

# ============================================================================
# SLIDE 17 — Thank You / Q&A
# ============================================================================
s = new_slide()
bg(s, NAVY)
fill_rect(s, 0, 0, PW, 0.16, BLUE)
fill_rect(s, 0, 7.34, PW, 0.16, TEAL)
tb(s, 1.0, 2.55, 11.3, 1.0, "Thank You", 54, WHITE, bold=True, font=FONT_SB)
fill_rect(s, 1.05, 3.75, 3.4, 0.04, BLUE)
tb(s, 1.0, 4.0, 11.3, 0.5, "Questions & Discussion are welcome", 18, ICE, italic=True)
tb(s, 1.0, 4.85, 11.3, 1.0,
   "Obi Kastanya  ·  Dhayu Intan Nareswari  ·  Ananta Dwi Prayoga Alwy  ·  Mahathir Muhammad",
   13.5, RGBColor(0xB9,0xC6,0xE2))
tb(s, 1.0, 5.25, 11.3, 0.5,
   "Department of Informatics — Institut Teknologi Sepuluh Nopember (ITS), Surabaya",
   12.5, RGBColor(0x9F,0xB4,0xDD))
for i, (lbl, col) in enumerate([("PREDICT", BLUE), ("EXPLAIN", TEAL),
                                ("PRESCRIBE", AMBER), ("NATURALIZE", PURPLE)]):
    chip(s, 1.0 + i*2.75, 6.2, 2.55, 0.5, lbl, col, size=12)
tb(s, 10.0, 7.04, 2.78, 0.3, f"{TOTAL} / {TOTAL}", 9, RGBColor(0x8A,0x9A,0xBE), align=PP_ALIGN.RIGHT)
notes(s,
"Demikian presentasi kami. Terima kasih atas perhatian Bapak/Ibu juri dan rekan-rekan "
"sekalian. Kami sangat terbuka untuk pertanyaan, masukan, maupun diskusi terkait metodologi, "
"hasil, maupun arah pengembangan ke depan.")

# ----------------------------------------------------------------------------
OUT = "kk-ngebut_slides_v2.pptx"
prs.save(OUT)
print(f"saved {OUT} with {len(prs.slides.__iter__().__class__ and list(prs.slides))} slides")
